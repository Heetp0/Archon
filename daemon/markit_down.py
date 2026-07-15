import logging
logger = logging.getLogger(__name__)

import os
import json
import hashlib
import shutil
import csv
from datetime import datetime

class MarkitDownNormalizer:
    def __init__(self, cache_dir: str):
        self.cache_dir = cache_dir
        self.index_path = os.path.join(cache_dir, 'index.json')
        os.makedirs(cache_dir, exist_ok=True)
        self._load_index()

    def _load_index(self):
        if os.path.exists(self.index_path):
            try:
                with open(self.index_path, 'r', encoding='utf-8') as f:
                    self.index = json.load(f)
            except Exception:
                self.index = {}
        else:
            self.index = {}

    def _save_index(self):
        try:
            with open(self.index_path, 'w', encoding='utf-8') as f:
                json.dump(self.index, f, indent=2)
        except Exception as e:
                    logger.warning(f"Handled exception: {e}")

    def _get_file_hash(self, filepath: str) -> str:
        sha256 = hashlib.sha256()
        try:
            with open(filepath, 'rb') as f:
                while chunk := f.read(8192):
                    sha256.update(chunk)
            return sha256.hexdigest()
        except Exception:
            return ''

    async def convert(self, source_path: str) -> str:
        """Converts document at source_path to Markdown and returns the path to the cached markdown file."""
        import asyncio
        if not os.path.exists(source_path):
            raise FileNotFoundError(f'Source file not found: {source_path}')

        file_hash = self._get_file_hash(source_path)
        if not file_hash:
            raise IOError(f'Failed to compute hash for file: {source_path}')

        # Check cache
        cache_filename = f'{file_hash}.md'
        cache_filepath = os.path.join(self.cache_dir, cache_filename)
        mtime = os.path.getmtime(source_path)

        if file_hash in self.index and os.path.exists(cache_filepath):
            # Cache hit
            return cache_filepath

        # Convert based on file type
        _, ext = os.path.splitext(source_path.lower())
        try:
            if ext == '.pdf':
                markdown_content = await asyncio.to_thread(self._convert_pdf, source_path)
            elif ext in ['.doc', '.docx']:
                markdown_content = self._convert_docx(source_path)
            elif ext in ['.ppt', '.pptx']:
                markdown_content = self._convert_pptx(source_path)
            elif ext in ['.xls', '.xlsx']:
                markdown_content = await asyncio.to_thread(self._convert_excel, source_path)
            elif ext in ['.csv']:
                markdown_content = self._convert_csv(source_path)
            elif ext in ['.html', '.htm']:
                markdown_content = self._convert_html(source_path)
            else:
                # Plain text or fallback
                markdown_content = self._convert_text(source_path)
        except Exception as e:
            # Fallback to plain text read if conversion library fails
            try:
                markdown_content = self._convert_text(source_path)
            except Exception:
                markdown_content = f'# File Conversion Failed\n\nError: {str(e)}\n'

        # Save to cache
        with open(cache_filepath, 'w', encoding='utf-8') as f:
            f.write(markdown_content)

        # Update index
        self.index[file_hash] = {
            'original_path': source_path,
            'cached_path': cache_filepath,
            'converted_at': datetime.now().isoformat(),
            'mtime': mtime
        }
        self._save_index()
        return cache_filepath

    def _convert_pdf(self, path: str) -> str:
        try:
            import pypdf
            reader = pypdf.PdfReader(path)
            text = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text.append(f'## Page {i+1}\n\n{page_text}')
            return '\n\n'.join(text)
        except ImportError:
            return '# PDF Reader Not Installed\n'

    def _convert_docx(self, path: str) -> str:
        try:
            import docx
            doc = docx.Document(path)
            text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            # Tables
            for table in doc.tables:
                table_md = []
                for row in table.rows:
                    row_text = [cell.text.strip().replace("|", "\\|").replace("\n", " ") for cell in row.cells]
                    table_md.append('| ' + ' | '.join(row_text) + ' |')
                if table_md:
                    # Add alignment row
                    cols = len(table.columns)
                    align = '| ' + ' | '.join(['---'] * cols) + ' |'
                    table_md.insert(1, align)
                    text.append('\n'.join(table_md))
            return '\n\n'.join(text)
        except ImportError:
            return '# Word Reader Not Installed\n'

    def _convert_pptx(self, path: str) -> str:
        try:
            import pptx
            prs = pptx.Presentation(path)
            text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    text.append(f'## Slide {i+1}\n\n' + '\n'.join(slide_text))
            return '\n\n'.join(text)
        except ImportError:
            return '# PPTX Reader Not Installed\n'

    def _convert_excel(self, path: str) -> str:
        try:
            import pandas as pd
            xls = pd.ExcelFile(path)
            sheets = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet_name)
                md = df.to_markdown(index=False)
                sheets.append(f'## Sheet: {sheet_name}\n\n{md}')
            return '\n\n'.join(sheets)
        except Exception:
            return '# Excel Reader Error\n'

    def _convert_csv(self, path: str) -> str:
        try:
            lines = []
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    row_escaped = [cell.replace("|", "\\|").replace("\n", " ") for cell in row]
                    lines.append('| ' + ' | '.join(row_escaped) + ' |')
            if lines:
                cols = len(lines[0].split('|')) - 2
                align = '| ' + ' | '.join(['---'] * cols) + ' |'
                lines.insert(1, align)
            return '\n'.join(lines)
        except Exception as e:
            return f'# CSV Read Error\n\n{str(e)}'

    def _convert_html(self, path: str) -> str:
        try:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            return soup.get_text()
        except Exception as e:
            return f'# HTML Read Error\n\n{str(e)}'

    def _convert_text(self, path: str) -> str:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        filename = os.path.basename(path)
        return f'# File: {filename}\n\n```\n{content}\n```'